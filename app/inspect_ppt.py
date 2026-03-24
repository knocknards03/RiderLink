import sys
try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not installed.")
    sys.exit(1)

TEMPLATE_PATH = "/Users/ashwinsmac/Downloads/Minor project Review 1 PPT-Template  B Tech (1).pptx"
prs = Presentation(TEMPLATE_PATH)

for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        print(f"[{shape.name}] {shape.text[:100]}...")
