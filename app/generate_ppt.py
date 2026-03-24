from pptx import Presentation

TEMPLATE_PATH = "/Users/ashwinsmac/Downloads/Minor project Review 1 PPT-Template  B Tech (1).pptx"
OUTPUT_PATH = "/Users/ashwinsmac/Desktop/Minor Project/RiderLink_Project_Review.pptx"

prs = Presentation(TEMPLATE_PATH)

def set_text(shape, new_text):
    if not shape.has_text_frame: return
    shape.text = new_text

# Slide 1: Title
title_shape = prs.slides[0].shapes[1]
set_text(title_shape, "RIDERLINK: OFFLINE MESH NETWORKING & TELEMETRY SYSTEM FOR MOTORCYCLES")
presented_shape = prs.slides[0].shapes[2]
set_text(presented_shape, "Presented By:\nAshwin & Team\n\nMinor project Review 1")

# Slide 3: Abstract
abstract_shape = prs.slides[2].shapes[1]
set_text(abstract_shape, "RiderLink is an advanced motorcycle telemetry and safety system leveraging ESP32-based LoRa mesh networking and Flutter. It addresses the critical issue of cellular black-spots during group rides by providing off-grid GPS tracking, automatic crash detection via IMU sensors, and peer-to-peer hazard broadcasting. The system includes an Android app for navigation and a digital garage for maintenance tracking.")

# Slide 4: Problem Statement
prob_shape = prs.slides[3].shapes[1]
set_text(prob_shape, "- Motorcycle riders frequently travel through remote areas lacking cellular coverage, making group coordination and emergency responses difficult.\n- Existing navigation apps rely heavily on constant active internet connections.\n- Current SOS systems are expensive, subscription-based, or lack automated vehicle-telemetry crash detection.")

# Slide 5: Introduction
intro_shape = prs.slides[4].shapes[1]
set_text(intro_shape, "Current Status: Base Flutter UI, SQLite Database, and off-grid LoRa mesh logic implemented.\nPlan: Finalize hardware integration with ESP32 sensors.\nScope: Real-time lean angle telemetry, Waze-style hazard reporting over LoRa, offline map routing, and group proximity alerts.\nAssumptions: Riders will equip their bikes with the RiderLink ESP32 hardware module.")

# Slide 6: Social/Environmental Impact
impact_shape = prs.slides[5].shapes[1]
set_text(impact_shape, "- Significantly reduces emergency response times for motorcycle accidents in remote areas via mesh-routed SOS beacons.\n- Encourages safer riding habits through live lean-angle and G-force telemetry feedback.\n- Promotes environmental awareness by integrating weather alerts and efficient curvy-route generation.")

# Slide 7: State of the Art
state_shape = prs.slides[6].shapes[1]
set_text(state_shape, "1. Traditional cellular GPS trackers (cannot operate fully off-grid).\n2. Garmin inReach (expensive satellite subscription required).\n3. Standard mesh intercoms (Cardo/Sena) lack map-visualized GPS coordinates.\nRiderLink bridges this gap by unifying free LoRa hardware with a rich smartphone UI.")

# Slide 8: Design
design_shape = prs.slides[7].shapes[1]
set_text(design_shape, "- Frontend: Flutter (Dart) for Android/iOS.\n- Backend/Hardware: ESP32 microcontroller with standard LoRa SX1278 transceiver.\n- Mapping: 'flutter_map' with offline Tile Layer support.\n- Database: SQLite (sqflite) for localized GPX route logging and maintenance tracking.\n- Protocol: Custom binary lightweight BLE mesh protocol.")

# Slide 10: Methodology
method_shape = prs.slides[9].shapes[1]
set_text(method_shape, "1. Hardware Assembly: ESP32 + LoRa + GPS + MPU6050 IMU.\n2. Firmware: C++ FreeRTOS tasks to handle high-frequency sensor fusion and LoRa interrupts.\n3. BLE Bridge: Flutter app acts as a central Bluetooth client, parsing custom characteristic payloads.\n4. UI Mapping: Layered state management (GetX) to render moving map markers and pop emergency Modals instantly.")

# Slide 11: References
ref_shape = prs.slides[10].shapes[1]
set_text(ref_shape, "[1] Augustin, A., et al., \"A study of LoRa: Long Range & Low Power Networks for the IoT,\" Sensors, 2016.\n[2] Flutter API Documentation, flutter.dev\n[3] Espressif ESP32 Datasheet and Hardware Reference Manual.")

prs.save(OUTPUT_PATH)
print("Successfully generated:", OUTPUT_PATH)
